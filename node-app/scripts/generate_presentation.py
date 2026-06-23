#!/usr/bin/env python3
"""
Generate a PowerPoint presentation documenting the COBOL to Node.js
migration process. Outputs to node-app/presentation/migration_overview.pptx.

Usage:
    python3 scripts/generate_presentation.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ── Helpers ──────────────────────────────────────────────────

def add_slide(prs, layout_index, title_text, body_items=None):
    """Add a slide with a title and optional bullet list."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    title = slide.shapes.title
    title.text = title_text

    if body_items and slide.placeholders[1]:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, item in enumerate(body_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.space_after = Pt(6)
            p.font.size = Pt(18)

    return slide


def add_title_slide(prs, title, subtitle):
    """Add a title slide (layout 0)."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, content_lines, notes_text=None):
    """Add a content slide with bullet points and optional speaker notes."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(content_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line.startswith("  "):
            p.text = line.strip()
            p.level = 1
            p.font.size = Pt(16)
        else:
            p.text = line
            p.level = 0
            p.font.size = Pt(18)
        p.space_after = Pt(4)

    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

    return slide


def add_code_slide(prs, title, code_text, notes_text=None):
    """Add a slide with a code block displayed in a text box."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Code box
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    code_tf = code_box.text_frame
    code_tf.word_wrap = True

    for i, line in enumerate(code_text.split("\n")):
        p = code_tf.paragraphs[0] if i == 0 else code_tf.add_paragraph()
        p.text = line
        p.font.name = "Courier New"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
        p.space_after = Pt(2)

    # Background fill for code area
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

    return slide


def add_table_slide(prs, title, headers, rows, notes_text=None):
    """Add a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    p = txBox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.3), Inches(9), Inches(0.4 * num_rows)
    )
    table = table_shape.table

    # Headers
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)

    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

    return slide


# ── Main ─────────────────────────────────────────────────────

def generate_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    add_title_slide(
        prs,
        "COBOL to Node.js Migration",
        "Modernizing a Legacy Accounting System with Devin"
    )

    # ── Slide 2: Agenda ──
    add_content_slide(prs, "Agenda", [
        "1. Why Modernize?",
        "2. COBOL Application Overview",
        "3. Migration Strategy",
        "4. Architecture Mapping",
        "5. Code Conversion Examples",
        "6. Testing Strategy",
        "7. Deployment Pipeline",
        "8. Results & Next Steps",
    ], notes_text="Walk through each phase of the migration process.")

    # ── Slide 3: Why Modernize ──
    add_content_slide(prs, "Why Modernize COBOL?", [
        "Aging workforce with COBOL expertise",
        "Difficulty integrating with modern APIs and cloud services",
        "Limited testing frameworks for COBOL",
        "High maintenance cost of legacy infrastructure",
        "Node.js advantages:",
        "  Vast ecosystem (npm) and community support",
        "  Easy cloud deployment (Docker, AWS, Azure)",
        "  Modern testing with Jest (coverage, mocking)",
        "  Async I/O and scalability",
    ], notes_text=(
        "COBOL is still running critical financial systems, but finding developers is "
        "increasingly difficult. Node.js provides a modern alternative with strong "
        "community support and cloud-native deployment options."
    ))

    # ── Slide 4: COBOL Application Overview ──
    add_table_slide(prs, "COBOL Application Overview", [
        "File", "Program ID", "Purpose", "Lines"
    ], [
        ["main.cob", "MainProgram", "CLI menu loop, user I/O", "36"],
        ["operations.cob", "Operations", "Business logic (credit/debit/view)", "40"],
        ["data.cob", "DataProgram", "In-memory balance storage", "23"],
    ], notes_text=(
        "The COBOL application is a simple 3-file accounting system. "
        "MainProgram handles the UI, Operations handles business logic, "
        "and DataProgram manages the balance state. Total: 99 lines of COBOL."
    ))

    # ── Slide 5: Migration Strategy ──
    add_content_slide(prs, "Migration Strategy: 5 Phases", [
        "Phase 1: Analyze COBOL structure and data types",
        "  Inventory files, map PIC types, trace call graph",
        "Phase 2: Create test plan BEFORE conversion",
        "  Document every business rule as a testable case",
        "Phase 3: Convert COBOL to Node.js module by module",
        "  Bottom-up: data.cob -> operations.cob -> main.cob",
        "Phase 4: Validate with unit and integration tests",
        "  Map each TESTPLAN.md case to a Jest test",
        "Phase 5: Deploy with Docker and CI/CD pipeline",
        "  Multi-environment: development, staging, production",
    ], notes_text=(
        "The key insight is to create the test plan BEFORE writing any Node.js code. "
        "This ensures no business logic is lost during conversion. We convert bottom-up "
        "so each layer can be tested independently."
    ))

    # ── Slide 6: Architecture Mapping ──
    add_table_slide(prs, "Architecture Mapping", [
        "COBOL Component", "Node.js Equivalent", "Pattern"
    ], [
        ["main.cob (MainProgram)", "src/main.js (AccountApp)", "Class with async run()"],
        ["operations.cob (Operations)", "src/operations.js (Operations)", "Dependency injection"],
        ["data.cob (DataProgram)", "src/data.js (DataStore)", "In-memory store"],
        ["CALL 'Program' USING", "method calls on injected deps", "DI replaces CALL"],
        ["PIC 9(6)V99", "number + toFixed(2)", "IEEE 754 float"],
        ["PERFORM UNTIL", "while (flag)", "Loop construct"],
        ["EVALUATE / WHEN", "switch / case", "Dispatch"],
    ], notes_text=(
        "Each COBOL program maps to a JavaScript class. "
        "CALL statements become method calls via dependency injection. "
        "COBOL's PIC type constraints are replaced by runtime validation."
    ))

    # ── Slide 7: Data Layer Conversion ──
    add_code_slide(prs, "Conversion: Data Layer", """COBOL (data.cob):
    01  STORAGE-BALANCE  PIC 9(6)V99 VALUE 1000.00.
    IF OPERATION-TYPE = 'READ'
        MOVE STORAGE-BALANCE TO BALANCE
    ELSE IF OPERATION-TYPE = 'WRITE'
        MOVE BALANCE TO STORAGE-BALANCE

Node.js (src/data.js):
    class DataStore {
      constructor(initialBalance = 1000.00) {
        this.balance = initialBalance;
      }
      read()            { return this.balance; }
      write(newBalance) { this.balance = newBalance; }
    }""", notes_text=(
        "The COBOL DataProgram uses string-based dispatch ('READ'/'WRITE'). "
        "In Node.js we use named methods instead. The VALUE clause becomes "
        "a constructor default parameter."
    ))

    # ── Slide 8: Business Logic Conversion ──
    add_code_slide(prs, "Conversion: Business Logic", """COBOL (operations.cob):
    CALL 'DataProgram' USING 'READ', FINAL-BALANCE
    IF FINAL-BALANCE >= AMOUNT
        SUBTRACT AMOUNT FROM FINAL-BALANCE
        CALL 'DataProgram' USING 'WRITE', FINAL-BALANCE
    ELSE
        DISPLAY "Insufficient funds for this debit."

Node.js (src/operations.js):
    debit(amount) {
      const currentBalance = this.dataStore.read();
      if (currentBalance >= parsedAmount) {
        const newBalance = currentBalance - parsedAmount;
        this.dataStore.write(newBalance);
        return { success: true, balance: newBalance };
      } else {
        return { success: false,
                 message: 'Insufficient funds for this debit.' };
      }
    }""", notes_text=(
        "The read-modify-write pattern is preserved. The key improvement is "
        "returning structured result objects instead of using DISPLAY directly. "
        "This separates business logic from I/O."
    ))

    # ── Slide 9: Testing Strategy ──
    add_table_slide(prs, "Testing Strategy", [
        "Test Case", "Description", "Type", "Source"
    ], [
        ["TC-1.1", "View balance displays 1000.00", "Unit + Integration", "operations.cob:17-18"],
        ["TC-2.1", "Credit adds to balance", "Unit + Integration", "operations.cob:23-26"],
        ["TC-2.2", "Zero credit, balance unchanged", "Unit", "operations.cob:23-26"],
        ["TC-3.1", "Valid debit subtracts", "Unit + Integration", "operations.cob:31-35"],
        ["TC-3.2", "Overdraft prevented", "Unit + Integration", "operations.cob:32,37"],
        ["TC-3.3", "Zero debit, balance unchanged", "Unit", "operations.cob:31-35"],
        ["TC-4.1", "Exit terminates cleanly", "Integration", "main.cob:30"],
    ], notes_text=(
        "Every test case was derived from the COBOL source BEFORE conversion. "
        "Unit tests validate individual operations. Integration tests simulate "
        "the full menu-driven workflow using stream-based I/O."
    ))

    # ── Slide 10: Deployment Pipeline ──
    add_content_slide(prs, "Deployment Pipeline", [
        "Multi-environment deployment script:",
        "  ./scripts/deploy.sh <env> <command>",
        "",
        "Environments:",
        "  development  - Direct Node.js, verbose logging",
        "  staging      - Docker container, mirrors production",
        "  production   - Docker with health checks, rollback",
        "",
        "Pipeline stages:",
        "  setup -> lint -> test -> build -> deploy -> health",
        "",
        "Cloud deployment support:",
        "  AWS ECS (via deploy/deploy.sh deploy-aws)",
        "  Azure App Service (via deploy/deploy.sh deploy-azure)",
    ], notes_text=(
        "The deployment script supports three environments with progressive "
        "hardening. Development runs locally, staging uses Docker to catch "
        "container-specific issues, and production adds health checks and "
        "rollback support."
    ))

    # ── Slide 11: Docker Architecture ──
    add_code_slide(prs, "Docker Architecture", """Multi-stage Dockerfile:

    Stage 1 (builder):
      FROM node:18-alpine
      COPY package*.json -> npm ci --only=production
      COPY src/

    Stage 2 (production):
      FROM node:18-alpine
      Create non-root user (appuser:appgroup)
      COPY from builder: node_modules, src, package.json
      COPY healthcheck.js
      EXPOSE 3000
      HEALTHCHECK: node healthcheck.js
      CMD: node src/main.js

    Docker Compose:
      ports: 3000:3000
      restart: unless-stopped
      resources: 128M memory, 0.5 CPU""", notes_text=(
        "Multi-stage build minimizes the final image size. "
        "Non-root user improves security. Health check endpoint "
        "at /health enables container orchestrator monitoring."
    ))

    # ── Slide 12: Results ──
    add_content_slide(prs, "Results", [
        "Successful 1:1 conversion of all business logic",
        "100% test coverage of COBOL test plan cases",
        "3 source files: data.js, operations.js, main.js",
        "22 unit tests + 7 integration tests",
        "Fully automated deployment pipeline",
        "Docker-ready for cloud deployment",
        "",
        "Modernization benefits achieved:",
        "  Testable: Jest with coverage reporting",
        "  Deployable: Docker + multi-env deploy script",
        "  Maintainable: Clean class-based architecture",
        "  Extensible: Ready for REST API, database, auth",
    ], notes_text=(
        "The conversion preserves 100% of the original business logic "
        "while adding modern development practices. The application is "
        "now testable, deployable to any cloud, and easy to extend."
    ))

    # ── Slide 13: Next Steps ──
    add_content_slide(prs, "Next Steps", [
        "1. Database persistence (PostgreSQL/MongoDB)",
        "  Replace in-memory DataStore with production DB",
        "2. REST API layer (Express.js)",
        "  Enable web and mobile clients",
        "3. Authentication (JWT tokens)",
        "  Secure multi-user access",
        "4. Transaction history and audit log",
        "  Compliance and traceability",
        "5. CI/CD with GitHub Actions",
        "  Automated testing and deployment on every push",
        "6. Monitoring and observability",
        "  Application metrics, log aggregation, alerting",
    ], notes_text=(
        "These enhancements build on the solid foundation created by "
        "the migration. Each can be implemented incrementally without "
        "disrupting existing functionality."
    ))

    # ── Slide 14: Thank You ──
    add_title_slide(
        prs,
        "Thank You",
        "Questions? See PROCESS.md and MIGRATION.md for full details."
    )

    # ── Save ──
    script_dir = os.path.dirname(os.path.abspath(__file__))
    project_dir = os.path.dirname(script_dir)
    output_dir = os.path.join(project_dir, "presentation")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "migration_overview.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    generate_presentation()
